#!/usr/bin/env python3
"""Deterministic extractor for tekshot knowledge_extract (file branch).

Reads ONE file and prints ONE JSON object on stdout - the contract is the
knowledgeExtraction struct in knowledge_extractor.go. Never raises: every
failure becomes {"ok": false, "error": <code>, "message": <vi sentence>}.

Libraries are the ones with musl wheels (markitdown is not): pdfplumber +
pypdfium2 for PDF (pypdfium2 renders scanned pages without poppler),
mammoth + markdownify for docx, openpyxl for xlsx, python-pptx for pptx.
"""
import argparse
import json
import os
import re
import sys

# A PDF page with fewer real characters than this is treated as a scan.
SCAN_TEXT_THRESHOLD = 30
# openpyxl read_only streams; a price list never needs more rows than this.
MAX_SHEET_ROWS = 2000


class ExtractError(Exception):
    def __init__(self, code, message):
        super().__init__(message)
        self.code = code
        self.message = message


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True)
    ap.add_argument("--mime", default="")
    ap.add_argument("--out-dir", required=True)
    ap.add_argument("--max-scan-pages", type=int, default=300)
    ap.add_argument("--dpi", type=int, default=150)
    a = ap.parse_args()
    try:
        result = extract(a.input, a.mime, a.out_dir, a.max_scan_pages, a.dpi)
    except ExtractError as e:
        result = {"ok": False, "error": e.code, "message": e.message}
    except Exception as e:  # noqa: BLE001 - contract: always JSON, never a traceback
        result = {"ok": False, "error": "corrupt",
                  "message": ("Không đọc được file: %s: %s" % (type(e).__name__, e))[:300]}
    # Bytes, not sys.stdout.write: the container has no LANG, so a Vietnamese
    # message would die with UnicodeEncodeError on an ASCII stdout.
    sys.stdout.buffer.write(json.dumps(result, ensure_ascii=False).encode("utf-8"))


def kind_for(mime, path):
    ext = os.path.splitext(path)[1].lower()
    mime = (mime or "").lower()
    if mime == "application/pdf" or ext == ".pdf":
        return "pdf"
    if mime.endswith("wordprocessingml.document") or ext == ".docx":
        return "docx"
    if mime.endswith("spreadsheetml.sheet") or ext == ".xlsx":
        return "xlsx"
    if mime.endswith("presentationml.presentation") or ext == ".pptx":
        return "pptx"
    if mime.startswith("image/") or ext in (".jpg", ".jpeg", ".png", ".webp"):
        return "image"
    if mime == "application/msword" or ext == ".doc":
        raise ExtractError("unsupported_mime",
                           "File .doc cũ không hỗ trợ — hãy lưu lại thành .docx rồi tải lên.")
    if mime == "application/vnd.ms-excel" or ext == ".xls":
        raise ExtractError("unsupported_mime",
                           "File .xls cũ không hỗ trợ — hãy lưu lại thành .xlsx rồi tải lên.")
    raise ExtractError("unsupported_mime",
                       "Định dạng không hỗ trợ. Nhận: ảnh, PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx).")


def extract(path, mime, out_dir, max_scan_pages, dpi):
    kind = kind_for(mime, path)
    stats = {"pages": 0, "text_pages": 0, "scan_pages": 0, "scan_pages_rendered": 0,
             "sheets": 0, "slides": 0}
    truncated, reason = False, ""
    if kind == "pdf":
        units, truncated, reason = extract_pdf(path, out_dir, max_scan_pages, dpi, stats)
    elif kind == "docx":
        units = extract_docx(path)
    elif kind == "xlsx":
        units, truncated, reason = extract_xlsx(path, stats)
    elif kind == "pptx":
        units = extract_pptx(path, stats)
    else:
        units = [{"kind": "image", "ref": os.path.basename(path), "image_path": path}]
    for i, u in enumerate(units):
        u["index"] = i
    return {"ok": True, "kind": kind, "units": units, "stats": stats,
            "truncated": truncated, "truncated_reason": reason}


# ---------- PDF ----------

def extract_pdf(path, out_dir, max_scan_pages, dpi, stats):
    import pdfplumber
    try:
        pdf = pdfplumber.open(path)
    except Exception as e:
        name = type(e).__name__
        if "Password" in name or "Encrypt" in name:
            raise ExtractError("encrypted_pdf", "PDF có mật khẩu — hãy gỡ mật khẩu rồi tải lại.")
        raise ExtractError("corrupt", "Không mở được PDF: %s" % name)
    units = []
    with pdf:
        stats["pages"] = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            text = page_markdown(page)
            if len(text.strip()) >= SCAN_TEXT_THRESHOLD:
                stats["text_pages"] += 1
                units.append({"kind": "text", "ref": "Trang %d" % (i + 1), "text": text})
            else:
                stats["scan_pages"] += 1
                units.append({"kind": "image", "ref": "Trang %d" % (i + 1),
                              "image_path": "", "_page": i})
    scan = [u for u in units if u["kind"] == "image"]
    truncated, reason = False, ""
    if scan:
        to_render = [u["_page"] for u in scan[:max_scan_pages]]
        if len(scan) > max_scan_pages:
            truncated = True
            reason = "PDF có %d trang scan, chỉ đọc %d trang đầu." % (len(scan), max_scan_pages)
        rendered = render_pages(path, to_render, out_dir, dpi)
        stats["scan_pages_rendered"] = len(rendered)
        kept = []
        for u in units:
            if u["kind"] != "image":
                kept.append(u)
                continue
            idx = u.pop("_page")
            if idx in rendered:
                u["image_path"] = rendered[idx]
                kept.append(u)
        units = kept
    return units, truncated, reason


def page_markdown(page):
    """Page text, then its tables re-emitted as Markdown so numbers keep their columns."""
    text = page.extract_text() or ""
    tables = []
    try:
        for table in page.extract_tables():
            md = table_to_markdown(table)
            if md:
                tables.append(md)
    except Exception:  # noqa: BLE001 - table detection is best-effort
        pass
    if tables:
        text = text + "\n\n" + "\n\n".join(tables)
    return text.strip()


def render_pages(path, indexes, out_dir, dpi):
    import pypdfium2 as pdfium
    try:
        pdf = pdfium.PdfDocument(path)
    except Exception as e:
        raise ExtractError("render_failed", "Không render được trang scan: %s" % type(e).__name__)
    out = {}
    for i in indexes:
        page = pdf[i]
        bitmap = page.render(scale=dpi / 72.0)
        img_path = os.path.join(out_dir, "page-%04d.png" % (i + 1))
        bitmap.to_pil().save(img_path, format="PNG")
        page.close()
        out[i] = img_path
    pdf.close()
    return out


# ---------- tables ----------

def table_to_markdown(rows):
    rows = [[_cell(c) for c in r] for r in (rows or []) if r and any(_cell(c) for c in r)]
    if len(rows) < 2:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    head = "| " + " | ".join(rows[0]) + " |"
    sep = "|" + "---|" * width
    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
    return "\n".join([head, sep, body])


def _cell(c):
    if c is None:
        return ""
    if isinstance(c, float) and c.is_integer():
        return str(int(c))
    return str(c).replace("\n", " ").replace("|", "/").strip()


# ---------- docx ----------

def extract_docx(path):
    import mammoth
    from markdownify import markdownify
    try:
        with open(path, "rb") as f:
            html = mammoth.convert_to_html(f).value
    except Exception as e:
        raise ExtractError("corrupt", "Không mở được Word: %s" % type(e).__name__)
    md = markdownify(html, heading_style="ATX", bullets="-")
    return split_markdown_by_heading(md, "Phan")


def split_markdown_by_heading(md, label):
    """One unit per H1/H2 section; no headings -> one unit, the Go chunker splits it further."""
    units, cur, cur_heading = [], [], ""
    for line in md.splitlines():
        m = re.match(r"^(#{1,2})\s+(.+)$", line)
        if m and "".join(cur).strip():
            units.append(_text_unit(cur, cur_heading, label, len(units)))
            cur = []
        if m:
            cur_heading = m.group(2).strip()
        cur.append(line)
    if "".join(cur).strip():
        units.append(_text_unit(cur, cur_heading, label, len(units)))
    return units


def _text_unit(lines, heading, label, n):
    u = {"kind": "text", "ref": "%s %d" % (label, n + 1), "text": "\n".join(lines).strip()}
    if heading:
        u["heading"] = heading
    return u


# ---------- xlsx ----------

def extract_xlsx(path, stats):
    from openpyxl import load_workbook
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        raise ExtractError("corrupt", "Không mở được Excel: %s" % type(e).__name__)
    units, truncated, reasons = [], False, []
    for ws in wb.worksheets:
        stats["sheets"] += 1
        rows = []
        for r in ws.iter_rows(values_only=True):
            row = [_cell(c) for c in (r or ())]
            while row and row[-1] == "":
                row.pop()
            if not any(row):
                continue
            rows.append(row)
            if len(rows) >= MAX_SHEET_ROWS:
                truncated = True
                reasons.append("Sheet %s dài quá %d dòng, chỉ lấy %d dòng đầu."
                               % (ws.title, MAX_SHEET_ROWS, MAX_SHEET_ROWS))
                break
        md = table_to_markdown(rows)
        if not md:
            continue
        units.append({"kind": "text", "ref": "Sheet %s" % ws.title, "heading": ws.title,
                      "text": "## %s\n\n%s" % (ws.title, md)})
    wb.close()
    return units, truncated, " ".join(reasons)


# ---------- pptx ----------

def extract_pptx(path, stats):
    from pptx import Presentation
    try:
        prs = Presentation(path)
    except Exception as e:
        raise ExtractError("corrupt", "Không mở được PowerPoint: %s" % type(e).__name__)
    units = []
    for n, slide in enumerate(prs.slides, start=1):
        stats["slides"] += 1
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                md = table_to_markdown([[c.text for c in row.cells] for row in shape.table.rows])
                if md:
                    parts.append(md)
            elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append("Ghi chú: " + notes)
        text = "\n\n".join(parts).strip()
        if not text:
            continue
        u = {"kind": "text", "ref": "Slide %d" % n, "text": "## Slide %d\n\n%s" % (n, text)}
        title = slide.shapes.title
        if title is not None and title.text.strip():
            u["heading"] = title.text.strip()
        units.append(u)
    return units


if __name__ == "__main__":
    main()
